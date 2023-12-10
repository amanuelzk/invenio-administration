import json
import requests
import re
import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from .model import *

def extract_text_from_pdf(file_path):
    reader = PdfReader(file_path)
    text = ''
    for page in reader.pages:
        text += page.extract_text() + '\n'
    return text

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
    return text

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def save_text_to_file(text, file_path):
    with open(file_path, 'w', encoding='utf-8') as file:
        file.write(text)

def process_downloaded_file(url, target_language):
    response = requests.get(url)
    if response.status_code == 200:
        file = extract_file(response)
        ext = os.path.splitext(file)[1].lower()
        if ext == '.pdf':
            text = extract_text_from_pdf(file)
        elif ext == '.docx':
            text = extract_text_from_docx(file)
        elif ext == '.pptx':
            text = extract_text_from_pptx(file)
        else:
            print(f"Unsupported file format: {file}")
            return

        output_file = os.path.splitext(file)[0] + f'_{target_language}.txt'
        save_text_to_file(text, output_file)
        print(f"Extracted text from {file} to {output_file}")

def extract_file(response):
    # Check if the header is present
    cd = response.headers.get("content-disposition")
    if not cd:
        return None

    # Find filename
    fname = re.findall('filename="(.+)"', cd)
    if len(fname) == 0:
        return None
    return fname[0]

def download_file_and_save(url, target_language, project_id):
    try:
        response = requests.get(url)
        # Check if the request was successful
        if response.status_code == 200:
            filename = extract_file(response)
            if filename:
                # Save the file
                translated_filename = f"{target_language}_{filename}"
                file_name_new = translated_filename.split('.')[0]
                file_extension = translated_filename.split('.')[-1]
                with open(translated_filename, "wb") as f:
                    f.write(response.content)
                if file_extension== 'json':
                    with open(translated_filename) as fp:
                        data = json.load(fp)
                
                print(f"File downloaded and saved as {translated_filename}")
                if target_language == 'ar-SA' and file_extension == 'docx':
                    update_file = update(ArabicFile).where(ArabicFile.project_id==project_id).values(file_name=file_name_new, file_data=response.content, translate_status='complete')
                elif target_language == 'en-US' and file_extension == 'docx':
                    update_file = update(EnglishFile).where(EnglishFile.project_id==project_id).values(file_name=file_name_new, file_data=response.content, translate_status='complete')
                elif target_language == 'fr-FR' and file_extension == 'docx':
                    update_file = update(FrenchFile).where(FrenchFile.project_id==project_id).values(file_name=file_name_new, file_data=response.content, translate_status='complete')
                elif target_language == 'es-ES' and file_extension == 'docx':
                    update_file = update(SpanishFile).where(SpanishFile.project_id==project_id).values(file_name=file_name_new, file_data=response.content, translate_status='complete')
                elif target_language == 'ar-SA' and file_extension == 'json':
                    update_file = update(ArabicMetadata).where(ArabicMetadata.project_id==project_id).values(file_name=file_name_new, file_data=data, translate_status='complete')
                elif target_language == 'en-US' and file_extension == 'json':
                    update_file = update(EnglishMetadata).where(EnglishMetadata.project_id==project_id).values(file_name=file_name_new, file_data=data, translate_status='complete')
                elif target_language == 'fr-FR' and file_extension == 'json':
                    update_file = update(FrenchMetadata).where(FrenchMetadata.project_id==project_id).values(file_name=file_name_new, file_data=data, translate_status='complete')
                elif target_language == 'es-ES' and file_extension == 'json':
                    update_file = update(SpanishMetadata).where(SpanishMetadata.project_id==project_id).values(file_name=file_name_new, file_data=data, translate_status='complete')

                session.execute(update_file)
                session.commit()
                
                return translated_filename
            else:
                print("Filename could not be extracted from the response.")
        else:
            print(f"Failed to download file. HTTP status code: {response.status_code}")
    except requests.RequestException as e:
        print(f"An error occurred while downloading the file: {e}")

def extract_and_save_text(filepath):
    file_extension = os.path.splitext(filepath)[1].lower()
    extracted_text = ""
    if file_extension == '.pdf':
        extracted_text = extract_text_from_pdf(filepath)
    elif file_extension == '.docx':
        extracted_text = extract_text_from_docx(filepath)
    elif file_extension == '.pptx':
        extracted_text = extract_text_from_pptx(filepath)
    elif file_extension == '.json':
        extracted_text = ""
    else:
        print(f"Unsupported file format: {filepath}")
    
    return extracted_text

def get_download_url(base_url, api_key, job_pass, job_id):
    headers = {
        'accept': 'application/json',
        'x-matecat-key': api_key,
    }
    status_url = f'{base_url}api/v2/jobs/{job_id}/{job_pass}'
    while True:
        response = requests.get(status_url, headers=headers)
        if response.status_code == 200:
            data = response.json()
            translation_download_url = data['job']['chunks'][0]['urls']['translation_download_url']
            return translation_download_url
        else:
            print(f"Failed to get job status for {job_id}: {response.status_code}")
            break

def get_job_status(base_url, api_key, status_data):
    headers = {
        'accept': 'application/json',
        'x-matecat-key': api_key,
    }
    response = requests.get(f'{base_url}api/status/', headers=headers, params=status_data)
    if response.status_code == 200:
        return response.json()
        
    else:
        print(f'Error: {response.status_code}')
        return None
    
# Main

# Get project status
def matecat_download(project_id, project_pass, uploaded,original_file_id):
    if uploaded == 'OK':
        downloaded = False

        status_data = {'id_project': project_id, 'project_pass': project_pass}
        response_data = get_job_status(base_url, api_key, status_data)
        # Extract project status
        job_status = response_data.get('status', 'Unknown')

        # Extract language pairs
        langpair_key = next(iter(response_data['jobs']['langpairs']), None)  # Gets the first key in the 'langpairs' dictionary
        langpair = response_data['jobs']['langpairs'].get(langpair_key) if langpair_key else 'Unknown'

        # Splitting to get target language
        if '|' in langpair:
            languages = langpair.split('|')
            target_language = languages[1] if len(languages) > 1 else 'Unknown'
        else:
            target_language = 'Unknown'
        if job_status == 'DONE':
            if downloaded is False:
                job_data = response_data.get('data')['jobs']
                for job_id , job_details in job_data.items():
                    job_pass = next(iter(job_details['totals']))
                    
                    download_translation_url = get_download_url(base_url, api_key, job_pass, job_id)
                    filename = download_file_and_save(download_translation_url, target_language, project_id)
                    file_extension= filename.split('.')[-1]
                    if file_extension != 'json':

                        searchable_text = extract_and_save_text(filename)
                        original_files=session.query(OriginalFile).filter_by(id=original_file_id)
                        for original_file in original_files:
                            update_file = update(OriginalFile).where(OriginalFile.id==original_file_id).values(searchability=func.concat(str(original_file.searchability),','+','.join([searchable_text])))
                        session.execute(update_file)
                        session.commit()
                    downloaded = True
            else:
                print("You have already downloaded the file")
        else:
            print(f"Error: {response_data['errors']}")
    else:
        print("Project is not uploaded yet.")


base_url = 'https://www.matecat.com/'
api_key = 'xY55uP6iraCKgrmErbQV-O4w4qe8cxXoCDnmW7Oal'

# Retrive id_project, project_pass and uploaded from database
# ..... your code here .....
def translate_function():
    arabic_records = session.query(ArabicFile).filter_by(translate_status=None)  
    french_records = session.query(FrenchFile).filter_by(translate_status=None)  
    spanish_records = session.query(SpanishFile).filter_by(translate_status=None)   
    english_records = session.query(EnglishFile).filter_by(translate_status=None)  

    arabic_metadatas = session.query(ArabicMetadata).filter_by(translate_status=None)  
    french_metadatas = session.query(FrenchMetadata).filter_by(translate_status=None)  
    spanish_metadatas = session.query(SpanishMetadata).filter_by(translate_status=None)   
    english_metadatas = session.query(EnglishMetadata).filter_by(translate_status=None)  


    for arabic_record in arabic_records:
        project_id = arabic_record.project_id
        project_pass = arabic_record.project_pass
        uploaded = arabic_record.project_status
        original_file_id = arabic_record.original_file_id
        matecat_download(project_id, project_pass, uploaded,original_file_id)

    for french_record in french_records:
        project_id = french_record.project_id
        project_pass = french_record.project_pass
        uploaded = french_record.project_status
        original_file_id = arabic_record.original_file_id
        matecat_download(project_id, project_pass, uploaded, original_file_id)

    for spanish_record in spanish_records:
        project_id = spanish_record.project_id
        project_pass = spanish_record.project_pass
        uploaded = spanish_record.project_status
        original_file_id = arabic_record.original_file_id
        matecat_download(project_id, project_pass, uploaded, original_file_id)

    for english_record in english_records:
        project_id = english_record.project_id
        project_pass = english_record.project_pass
        uploaded = english_record.project_status
        matecat_download(project_id, project_pass, uploaded, original_file_id)

    for arabic_metadata in arabic_metadatas:
        project_id = arabic_metadata.project_id
        project_pass = arabic_metadata.project_pass
        uploaded = arabic_metadata.project_status
        original_file_id = arabic_record.original_file_id
        matecat_download(project_id, project_pass, uploaded, original_file_id)

    for french_metadata in french_metadatas:
        project_id = french_metadata.project_id
        project_pass = french_metadata.project_pass
        uploaded = french_metadata.project_status
        original_file_id = arabic_record.original_file_id
        matecat_download(project_id, project_pass, uploaded, original_file_id)

    for spanish_metadata in spanish_metadatas:
        project_id = spanish_metadata.project_id
        project_pass = spanish_metadata.project_pass
        uploaded = spanish_metadata.project_status
        original_file_id = arabic_record.original_file_id
        matecat_download(project_id, project_pass, uploaded, original_file_id)

    for english_metadata in english_metadatas:
        project_id = english_metadata.project_id
        project_pass = english_metadata.project_pass
        uploaded = english_metadata.project_status
        original_file_id = arabic_record.original_file_id
        matecat_download(project_id, project_pass, uploaded, original_file_id)


